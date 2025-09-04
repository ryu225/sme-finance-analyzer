import os
from io import BytesIO
from datetime import datetime

import numpy as np
import pandas as pd
import streamlit as st
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
from plotly import colors as pcolors


#page setup

st.set_page_config(page_title="SME Finance Analyzer", layout="wide")
st.title("📊 SME Finance Analyzer")
st.caption("Upload transactions, review monthly P&L, explore trends, and export a concise PowerPoint.")

REQUIRED_COLS = ["date", "type", "category", "description", "amount"]
OPTIONAL_COLS = ["tax", "currency"]

def load_csv(path: str, fallback: pd.DataFrame | None = None) -> pd.DataFrame:
    if os.path.exists(path):
        return pd.read_csv(path)
    return fallback if fallback is not None else pd.DataFrame()

def load_sample() -> pd.DataFrame:
    sample_path = "data/sample_transactions.csv"
    if os.path.exists(sample_path):
        df = pd.read_csv(sample_path)
    else:
        # Fallback if file is missing
        df = pd.DataFrame({
            "date": ["2025-01-03","2025-01-05","2025-01-12"],
            "type": ["sale","expense","sale"],
            "category": ["online order","software","retail store"],
            "description": ["Order #1001","SaaS subscription","POS ticket #1102"],
            "amount": [950.0,-120.0,1250.0],
            "tax": [45.0,0.0,62.5],
            "currency": ["SAR","SAR","SAR"],
        })
    return df

def clean_df(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    # normalize columns
    df.columns = [c.strip().lower() for c in df.columns]
    # ensure required columns exist
    for col in REQUIRED_COLS:
        if col not in df.columns:
            raise ValueError(f"Missing required column: {col}")
    # parse date
    df["date"] = pd.to_datetime(df["date"], errors="coerce")
    df = df.dropna(subset=["date"])
    # numeric amounts
    df["amount"] = pd.to_numeric(df["amount"], errors="coerce")
    df = df.dropna(subset=["amount"])
    # type normalization
    df["type"] = df["type"].str.lower().str.strip()
    df["type"] = df["type"].replace({"sales":"sale","income":"sale","exp":"expense","expenses":"expense"})
    # defaults
    if "tax" not in df.columns:
        df["tax"] = 0.0
    if "currency" not in df.columns:
        df["currency"] = "SAR"
    df["currency"] = df["currency"].astype(str).str.upper().str.strip()
    # month string
    df["month"] = df["date"].dt.to_period("M").astype(str)
    return df

def apply_currency(df: pd.DataFrame, fx: pd.DataFrame, base_currency: str) -> pd.DataFrame:
    fx = fx.copy()
    if fx.empty:
        # assume all already in base
        df["amount_base"] = df["amount"]
        return df
    fx.columns = ["currency","rate_to_base"]
    fx["currency"] = fx["currency"].str.upper().str.strip()
    # ensure base currency exists
    if base_currency not in fx["currency"].values:
        fx = pd.concat([fx, pd.DataFrame({"currency":[base_currency],"rate_to_base":[1.0]})], ignore_index=True)
    df = df.merge(fx, how="left", left_on="currency", right_on="currency")
    df["rate_to_base"] = df["rate_to_base"].fillna(1.0)
    df["amount_base"] = df["amount"] * df["rate_to_base"]
    return df

def apply_vat(df: pd.DataFrame, include_vat: bool, vat_rate: float) -> pd.DataFrame:
    df = df.copy()
    vr = vat_rate / 100.0
    if include_vat:
        # amounts already gross → compute net
        df["net_amount"] = np.where(df["type"]=="sale", df["amount_base"]/(1+vr), df["amount_base"]/(1+vr))
        df["vat_amount"] = df["amount_base"] - df["net_amount"]
    else:
        # amounts are net → compute gross for reference
        df["net_amount"] = df["amount_base"]
        df["vat_amount"] = np.where(df["type"]=="sale", df["net_amount"]*vr, df["net_amount"]*vr)
    return df

def summarize(df: pd.DataFrame) -> dict:
    rev = df.loc[df["type"]=="sale","net_amount"].sum()
    exp = -df.loc[df["type"]=="expense","net_amount"].sum()
    profit = rev - exp
    margin = (profit/rev*100.0) if rev != 0 else np.nan
    return dict(revenue=rev, expenses=exp, profit=profit, margin=margin)

def monthly_pnl(df: pd.DataFrame) -> pd.DataFrame:
    rev = df[df["type"]=="sale"].groupby("month")["net_amount"].sum().rename("revenue")
    exp = -df[df["type"]=="expense"].groupby("month")["net_amount"].sum().rename("expenses")
    pnl = pd.concat([rev, exp], axis=1).fillna(0.0)
    pnl["profit"] = pnl["revenue"] - pnl["expenses"]
    pnl["margin_%"] = np.where(pnl["revenue"]!=0, pnl["profit"]/pnl["revenue"]*100, np.nan)
    pnl = pnl.reset_index()
    return pnl

def expense_breakdown(df: pd.DataFrame) -> pd.DataFrame:
    e = df[df["type"]=="expense"].copy()
    if e.empty:
        return pd.DataFrame({"category":[], "total":[]})
    out = (-e.groupby("category")["net_amount"].sum()).reset_index().rename(columns={"net_amount":"total"})
    return out.sort_values("total", ascending=False)

def anomalies(df: pd.DataFrame, z_thresh: float = 2.5) -> pd.DataFrame:
    e = df[df["type"]=="expense"].copy()
    if e.empty or e["net_amount"].std(ddof=0) == 0:
        return pd.DataFrame(columns=e.columns)
    z = (e["net_amount"] - e["net_amount"].mean()) / e["net_amount"].std(ddof=0)
    # large (negative) expenses → very low z
    mask = z < -abs(z_thresh)
    out = e.loc[mask].copy()
    out["z_score"] = z[mask].round(2)
    return out.sort_values("z_score")

def fig_to_png_bytes(fig) -> bytes:
    # requires 'kaleido'
    buf = BytesIO()
    fig.write_image(buf, format="png", scale=2)
    return buf.getvalue()

def build_ppt(pnl_df: pd.DataFrame, kpis: dict, fig_rev_exp, fig_profit, exp_cat: pd.DataFrame) -> BytesIO:

    # Export style
    palette = pcolors.qualitative.Set2
    template = "plotly_white"

    # Rebuild export figures
    fig1 = px.bar(
        pnl_df, x="month", y=["revenue", "expenses"],
        barmode="group", title="Revenue and Expenses by Month",
        template=template, color_discrete_sequence=palette
    )
    fig2 = px.bar(
        pnl_df, x="month", y="profit",
        title="Monthly Profit", template=template,
        color_discrete_sequence=palette
    )
    fig3 = None
    if not exp_cat.empty:
        fig3 = px.pie(
            exp_cat, names="category", values="total",
            title=None, template=template,
            color_discrete_sequence=palette
        )

    def fig_to_png_bytes(fig) -> bytes:
        buf = BytesIO()
        fig.write_image(buf, format="png", scale=2)  # requires kaleido
        return buf.getvalue()

    prs = Presentation()

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "SME Finance Analyzer"
    slide.placeholders[1].text = f"Auto-generated on {datetime.now():%Y-%m-%d}"

    # KPIs
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Metrics"
    tf = slide2.placeholders[1].text_frame
    tf.clear()
    for line in [
        f"Revenue: {kpis['revenue']:,.0f}",
        f"Expenses: {kpis['expenses']:,.0f}",
        f"Profit: {kpis['profit']:,.0f}",
        f"Margin: {kpis['margin']:.1f}%" if not np.isnan(kpis['margin']) else "Margin: n/a",
    ]:
        p = tf.add_paragraph(); p.text = line; p.level = 0

    # Monthly charts
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Monthly Trends"
    slide3.shapes.add_picture(BytesIO(fig_to_png_bytes(fig1)), Inches(0.5), Inches(1.5), width=Inches(4.6), height=Inches(3.1))
    slide3.shapes.add_picture(BytesIO(fig_to_png_bytes(fig2)), Inches(5.0), Inches(1.5), width=Inches(4.6), height=Inches(3.1))

    # Expense breakdown (pie)
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "Top Expense Categories"
    if fig3 is not None:
        slide4.shapes.add_picture(BytesIO(fig_to_png_bytes(fig3)), Inches(2.0), Inches(1.6), width=Inches(6.5))

    # P&L table
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    slide5.shapes.title.text = "P&L by Month"
    rows, cols = len(pnl_df) + 1, 4
    tbl = slide5.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8 + 0.28*rows)).table
    headers = ["Month", "Revenue", "Expenses", "Profit"]
    for j, htxt in enumerate(headers):
        tbl.cell(0, j).text = htxt
    for i, (_, r) in enumerate(pnl_df.iterrows(), start=1):
        tbl.cell(i, 0).text = str(r["month"])
        tbl.cell(i, 1).text = f"{r['revenue']:,.0f}"
        tbl.cell(i, 2).text = f"{r['expenses']:,.0f}"
        tbl.cell(i, 3).text = f"{r['profit']:,.0f}"

    out = BytesIO()
    prs.save(out); out.seek(0)
    return out

#  Data input
st.sidebar.subheader("Data")
uploaded = st.sidebar.file_uploader("Upload CSV or Excel", type=["csv","xlsx"])
fx_df = load_csv("data/fx_rates.csv")
budget_df = load_csv("data/budget.csv")
base_currency = st.sidebar.text_input("Base currency (for FX conversion)", value="SAR").upper().strip()

st.sidebar.subheader("VAT")
include_vat = st.sidebar.checkbox("Amounts include VAT", value=True)
vat_rate = st.sidebar.number_input("VAT %", 0.0, 25.0, 15.0, 0.5)

st.sidebar.subheader("Tools")
z_thresh = st.sidebar.slider("Anomaly sensitivity (z-score)", 0.5, 4.0, 2.5, 0.1)

with st.sidebar.expander("Download templates"):
    st.download_button("Transactions CSV template", data=load_csv("data/template_transactions.csv").to_csv(index=False).encode("utf-8"), file_name="template_transactions.csv")
    st.download_button("FX rates CSV", data=fx_df.to_csv(index=False).encode("utf-8") if not fx_df.empty else b"currency,rate_to_base\nSAR,1.0", file_name="fx_rates.csv")
    st.download_button("Budget CSV", data=budget_df.to_csv(index=False).encode("utf-8") if not budget_df.empty else b"month,revenue_target,expense_target\n", file_name="budget.csv")

#load transactions
if uploaded:
    if uploaded.name.lower().endswith(".xlsx"):
        df = pd.read_excel(uploaded)
    else:
        df = pd.read_csv(uploaded)
else:
    df = load_sample()

#cleaning
try:
    df = clean_df(df)
except Exception as e:
    st.error(f"Data error: {e}")
    st.stop()

df = apply_currency(df, fx_df, base_currency)
df = apply_vat(df, include_vat, vat_rate)

#  KPIs
kpis = summarize(df)
c1, c2, c3, c4 = st.columns(4)
c1.metric(f"Revenue ({base_currency})", f"{kpis['revenue']:,.0f}")
c2.metric(f"Expenses ({base_currency})", f"{kpis['expenses']:,.0f}")
c3.metric(f"Profit ({base_currency})", f"{kpis['profit']:,.0f}")
c4.metric("Margin", f"{kpis['margin']:.1f}%" if not np.isnan(kpis['margin']) else "n/a")

#  Visuals
pnl = monthly_pnl(df)
exp_cat = expense_breakdown(df)

# Bugetig (opt)
if not budget_df.empty:
    b = budget_df.copy()
    b["month"] = b["month"].astype(str)
    pnl = pnl.merge(b, how="left", on="month")
    pnl["rev_variance"] = pnl["revenue"] - pnl.get("revenue_target", 0)
    pnl["exp_variance"] = pnl.get("expense_target", 0) - pnl["expenses"]

st.subheader("Monthly Revenue vs Expenses")
fig_rev_exp = px.bar(pnl, x="month", y=["revenue","expenses"], barmode="group", title="Revenue and Expenses by Month")
st.plotly_chart(fig_rev_exp, use_container_width=True)

st.subheader("Profit by Month")
fig_profit = px.bar(pnl, x="month", y="profit", title="Monthly Profit")
st.plotly_chart(fig_profit, use_container_width=True)

if "revenue_target" in pnl.columns:
    st.info("Budget loaded: showing variance columns in the P&L table below.")

#tables
with st.expander("P&L Table"):
    st.dataframe(pnl)

st.subheader("Expense Breakdown by Category")
if not exp_cat.empty:
    fig_exp = px.pie(exp_cat, names="category", values="total", title=None)
    st.plotly_chart(fig_exp, use_container_width=True)
else:
    st.info("No expenses found in the dataset.")

# Anomalies
st.subheader("Anomaly Detection (Expenses)")
anom = anomalies(df, z_thresh=z_thresh)
if not anom.empty:
    st.dataframe(anom[["date","category","description","net_amount","currency","z_score"]].sort_values("z_score"))
    st.caption("Rows flagged as unusually large expenses relative to the dataset mean/std.")
else:
    st.caption("No anomalies detected with current threshold.")


#  Export 
st.subheader("Export")
col_a, col_b, col_c = st.columns(3)

with col_a:
    st.download_button(
        "Download Clean CSV",
        data=df.to_csv(index=False).encode("utf-8"),
        file_name=f"transactions_clean_{datetime.now():%Y%m%d}.csv",
        mime="text/csv",
    )

with col_b:
    try:
        ppt_bytes = build_ppt(pnl, kpis, fig_rev_exp, fig_profit, exp_cat)
        st.download_button(
            "Export PowerPoint",
            data=ppt_bytes,
            file_name=f"SME_Finance_Analyzer_{datetime.now():%Y%m%d}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as e:
        st.error(f"PPT export failed: {type(e).__name__}: {e}")
        st.caption("Tip: ensure 'kaleido' and 'python-pptx' are installed in the SAME environment used by Streamlit.")

with col_c:
    st.download_button(
        "Download P&L Table (CSV)",
        data=pnl.to_csv(index=False).encode("utf-8"),
        file_name=f"pnl_{datetime.now():%Y%m%d}.csv",
        mime="text/csv",
    )

#  Raw table 
with st.expander("Preview raw transactions"):
    st.dataframe(df.head(500))

